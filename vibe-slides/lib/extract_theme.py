#!/usr/bin/env python3
"""
extract_theme.py

Mechanical extraction of theme-relevant information from a .pptx template.

Emits two files into the output directory:

  theme.raw.json         — exhaustive structured dump of everything readable
  extraction_report.md   — judgment calls the agent must resolve with the user

This script DOES NOT make aesthetic decisions. It reports everything it
finds, flags ambiguities, and proposes mappings — but the agent + user
produce the final theme.json.

Usage:
    python3 extract_theme.py --input template.pptx --output-dir <project-dir>

Requires: python-pptx
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
import zipfile
from collections import Counter
from pathlib import Path
from xml.etree import ElementTree as ET

NS = {
    "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p":   "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r":   "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
EMU = 914400


# ---------- low-level helpers -----------------------------------------------


def emu_in(v):
    """EMU string → inches (3 decimals), or None."""
    try:
        return round(int(v) / EMU, 3)
    except (TypeError, ValueError):
        return None


def resolve_rel(base_path: str, target: str) -> str:
    """Resolve a rels Target against the base file's location."""
    if target.startswith("/"):
        return target.lstrip("/")
    base_dir = "/".join(base_path.split("/")[:-1])
    return os.path.normpath(os.path.join(base_dir, target)).replace(os.sep, "/")


def rels_path_for(xml_path: str) -> str:
    """ppt/foo/bar.xml → ppt/foo/_rels/bar.xml.rels"""
    parts = xml_path.split("/")
    return "/".join(parts[:-1] + ["_rels", parts[-1] + ".rels"])


def parse_rels(z: zipfile.ZipFile, rels_path: str) -> list[dict]:
    try:
        with z.open(rels_path) as fh:
            root = ET.fromstring(fh.read())
    except KeyError:
        return []
    return [
        {"id": r.get("Id"), "type": r.get("Type"), "target": r.get("Target")}
        for r in root.findall("rel:Relationship", NS)
    ]


def layout_num(path: str) -> int:
    m = re.search(r"slideLayout(\d+)\.xml$", path)
    return int(m.group(1)) if m else 0


# ---------- core extraction -------------------------------------------------


def find_active_theme(z: zipfile.ZipFile):
    """
    Follow presentation.xml → first slideMaster → its theme.
    Returns (master_path, theme_path, master_count).
    """
    pres_rels = parse_rels(z, "ppt/_rels/presentation.xml.rels")
    masters = [r for r in pres_rels if r["type"].endswith("/slideMaster")]
    if not masters:
        raise RuntimeError("no slideMaster relationships in ppt/_rels/presentation.xml.rels")
    first = masters[0]
    master_path = resolve_rel("ppt/presentation.xml", first["target"])
    m_rels = parse_rels(z, rels_path_for(master_path))
    theme_rels = [r for r in m_rels if r["type"].endswith("/theme")]
    if not theme_rels:
        raise RuntimeError(f"master {master_path} has no theme rel")
    theme_path = resolve_rel(master_path, theme_rels[0]["target"])
    return master_path, theme_path, len(masters)


def parse_theme(z: zipfile.ZipFile, theme_path: str) -> dict:
    """Extract colors and declared fonts from a theme XML."""
    with z.open(theme_path) as fh:
        root = ET.fromstring(fh.read())
    out = {"colors": {}, "fonts": {}}

    clr = root.find(".//a:clrScheme", NS)
    if clr is not None:
        for child in clr:
            slot = child.tag.split("}")[-1]
            if not len(child):
                continue
            inner = list(child)[0]
            if inner.tag.endswith("srgbClr"):
                out["colors"][slot] = inner.get("val")
            elif inner.tag.endswith("sysClr"):
                out["colors"][slot] = inner.get("lastClr") or inner.get("val")

    fs = root.find(".//a:fontScheme", NS)
    if fs is not None:
        for key, xpath in [
            ("major",   "a:majorFont/a:latin"),
            ("minor",   "a:minorFont/a:latin"),
            ("majorEA", "a:majorFont/a:ea"),
            ("minorEA", "a:minorFont/a:ea"),
        ]:
            el = fs.find(xpath, NS)
            if el is not None and el.get("typeface"):
                out["fonts"][key] = el.get("typeface")
    return out


def scan_layout(z: zipfile.ZipFile, layout_path: str) -> dict:
    """Pull name, shapes (placeholders + decorations), and declared font sizes."""
    with z.open(layout_path) as fh:
        root = ET.fromstring(fh.read())

    cSld = root.find("p:cSld", NS)
    name = cSld.get("name") if cSld is not None else ""

    shapes = []
    sz_values: set[float] = set()

    for sp in root.findall(".//p:sp", NS):
        info: dict = {}
        cnv = sp.find("p:nvSpPr/p:cNvPr", NS)
        info["name"] = cnv.get("name") if cnv is not None else "?"

        ph = sp.find("p:nvSpPr/p:nvPr/p:ph", NS)
        if ph is not None:
            info["kind"]   = "placeholder"
            info["phType"] = ph.get("type") or "body"
            info["phIdx"]  = ph.get("idx")
        else:
            info["kind"] = "shape"
            geo = sp.find("p:spPr/a:prstGeom", NS)
            if geo is not None:
                info["geom"] = geo.get("prst")

        xfrm = sp.find("p:spPr/a:xfrm", NS)
        if xfrm is not None:
            off = xfrm.find("a:off", NS)
            ext = xfrm.find("a:ext", NS)
            if off is not None:
                info["x"] = emu_in(off.get("x"))
                info["y"] = emu_in(off.get("y"))
            if ext is not None:
                info["w"] = emu_in(ext.get("cx"))
                info["h"] = emu_in(ext.get("cy"))
            rot = xfrm.get("rot")
            if rot is not None:
                info["rotate"] = round(int(rot) / 60000, 1)

        fill = sp.find("p:spPr/a:solidFill", NS)
        if fill is not None:
            srgb = fill.find("a:srgbClr", NS)
            scheme = fill.find("a:schemeClr", NS)
            if srgb is not None:
                info["fill"] = {"type": "srgb", "val": srgb.get("val")}
            elif scheme is not None:
                info["fill"] = {"type": "scheme", "val": scheme.get("val")}

        for rPr in sp.findall(".//a:defRPr", NS):
            if rPr.get("sz"):
                try:
                    sz_values.add(int(rPr.get("sz")) / 100)
                except ValueError:
                    pass

        shapes.append(info)

    return {
        "name":     name,
        "path":     layout_path,
        "shapes":   shapes,
        "szValues": sorted(sz_values),
    }


def slide_dims(z: zipfile.ZipFile):
    with z.open("ppt/presentation.xml") as fh:
        root = ET.fromstring(fh.read())
    sz = root.find("p:sldSz", NS)
    return emu_in(sz.get("cx")), emu_in(sz.get("cy"))


def scan_observed_fonts(pptx_path: str) -> dict:
    """Count run-level font names across all slides (observed, not declared)."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    counter: Counter[str] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name:
                        counter[r.font.name] += 1
    return dict(counter.most_common())


# ---------- report ----------------------------------------------------------


def build_report(raw: dict, out_dir: Path) -> None:
    L = []  # lines
    L.append("# Theme extraction report")
    L.append("")
    L.append(f"**Source:** `{raw['source']}`")
    L.append(f"**Slide:** {raw['slide']['w']} × {raw['slide']['h']} in")
    L.append(f"**Master:** `{raw['master']}`")
    L.append(f"**Theme:** `{raw['theme']}`")
    if raw["masterCount"] > 1:
        L.append("")
        L.append(f"⚠️ **{raw['masterCount']} slide masters found.** Extractor used the first. "
                 "If this deck mixes layouts from different masters, re-run with the right one "
                 "specified manually.")

    L.append("")
    L.append("---")
    L.append("")
    L.append("## Judgment calls — agent must resolve with the user")

    # ---- 1. Font ----------------------------------------------------------
    L.append("")
    L.append("### 1. Font: declared vs. observed")
    declared = raw["fonts"]["declared"]
    observed = raw["fonts"]["observed"]
    L.append(f"- **Theme declares:** major=`{declared.get('major', '?')}`, "
             f"minor=`{declared.get('minor', '?')}`"
             + (f", majorEA=`{declared['majorEA']}`" if declared.get("majorEA") else "")
             + (f", minorEA=`{declared['minorEA']}`" if declared.get("minorEA") else ""))
    if observed:
        L.append(f"- **Top fonts actually typed in slides:**")
        for name, count in list(observed.items())[:5]:
            L.append(f"    - `{name}` — {count} runs")
    else:
        L.append("- No fonts observed on any slide.")

    top_obs = next(iter(observed.items()), (None, 0))[0]
    decl_latin = {declared.get("major"), declared.get("minor")}
    if top_obs and top_obs not in decl_latin:
        L.append("")
        L.append(f"⚠️ **Mismatch.** Author's actual typed font (`{top_obs}`) differs from "
                 f"theme declaration (`{declared.get('minor', '?')}`).")
        L.append(f"**Recommend:** commit to `{top_obs}` (what they actually typed). "
                 "Ask user to confirm.")
    else:
        L.append("")
        L.append("No mismatch. Can defer to theme declaration.")

    # ---- 2. Motif candidates --------------------------------------------
    L.append("")
    L.append("### 2. Motif candidates")
    L.append("")
    L.append("Decorative shapes found in layouts (non-placeholder shapes with fills). "
             "Shapes that **bleed off the slide edge** (negative x/y) are especially "
             "likely to be signature motifs.")

    total_decos = 0
    for layout in raw["layouts"]:
        decos = [s for s in layout["shapes"]
                 if s.get("kind") == "shape" and "fill" in s]
        if not decos:
            continue
        L.append("")
        L.append(f"**Layout `{layout['name']}`** ({layout['path'].split('/')[-1]}):")
        for shp in decos:
            bleeds = []
            if shp.get("x") is not None and shp["x"] < 0: bleeds.append("left")
            if shp.get("y") is not None and shp["y"] < 0: bleeds.append("top")
            bleed_tag = f" **[bleeds {'+'.join(bleeds)}]**" if bleeds else ""
            pos = f"({shp.get('x')}, {shp.get('y')}, {shp.get('w')}×{shp.get('h')})"
            fill = shp.get("fill", {})
            fill_str = f"fill=`{fill.get('type')}:{fill.get('val')}`"
            rot = f" rot={shp['rotate']}°" if shp.get("rotate") else ""
            L.append(f"  - `{shp['name']}` {shp.get('geom', '?')} @ {pos} {fill_str}{rot}{bleed_tag}")
            total_decos += 1

    if total_decos == 0:
        L.append("")
        L.append("None found. The template has no decorative shapes in its layouts.")
    else:
        L.append("")
        L.append("**For each visually-cohesive group, ask the user:**")
        L.append("- Promote to a reusable `theme.motif.<name>` flag?")
        L.append("- Suggested names: `bookmark`, `cornerMark`, `accentBar`, `stripe`, `flag`.")
        L.append("- Which layouts should render the motif? (e.g., content-only vs. all)")

    # ---- 3. Color semantics ---------------------------------------------
    L.append("")
    L.append("### 3. Color semantics")
    L.append("")
    L.append("**Full palette extracted from theme:**")
    for slot, hex_ in raw["colors"].items():
        L.append(f"- `{slot}` = `#{hex_}`")

    usage = Counter()
    for layout in raw["layouts"]:
        for shp in layout["shapes"]:
            f = shp.get("fill", {})
            if f.get("type") == "scheme":
                usage[f["val"]] += 1
    if usage:
        L.append("")
        L.append("**Scheme color usage across layouts** (fill references only — "
                 "text color usage requires deeper XML inspection):")
        for name, count in usage.most_common():
            L.append(f"- `{name}` used {count}×")

    L.append("")
    L.append("**Recommended initial mapping (user must confirm or adjust):**")
    colors = raw["colors"]
    primary = colors.get("accent1") or colors.get("dk2")
    rows = [
        ("primary",    primary),
        ("accent",     primary),
        ("text",       colors.get("dk1", "000000")),
        ("muted",      colors.get("dk2", "6B7280")),
        ("bg",         colors.get("lt1", "FFFFFF")),
        ("surface",    colors.get("lt1", "FFFFFF")),
        ("secondary",  colors.get("lt2", "F0F0F0")),
        ("coral",      colors.get("accent5")),
        ("teal",       colors.get("accent3")),
        ("onPrimary",  colors.get("lt1", "FFFFFF")),
        ("onAccent",   colors.get("lt1", "FFFFFF")),
    ]
    for name, val in rows:
        if val:
            L.append(f"- `{name}` = `#{val}`")
    L.append("")
    L.append("Single-accent themes (where accent1 plays both primary and accent "
             "roles) are common. If the palette has visually distinct roles, split "
             "them — ask the user.")

    # ---- 4. Size hierarchy ----------------------------------------------
    L.append("")
    L.append("### 4. Font size hierarchy")
    all_sz = sorted({sz for layout in raw["layouts"] for sz in layout["szValues"]},
                    reverse=True)
    L.append(f"Declared font sizes across layouts: `{all_sz}` pt")
    if all_sz:
        labels = ["title", "section", "h3", "body", "small", "caption"]
        L.append("")
        L.append("Suggested mapping (largest → smallest):")
        for label, sz in zip(labels, all_sz):
            L.append(f"- `size.{label}` = {sz}pt")
        if len(all_sz) > len(labels):
            L.append(f"- (extra sizes not mapped: {all_sz[len(labels):]})")

    # ---- 5. Spacing ------------------------------------------------------
    L.append("")
    L.append("### 5. Spacing from layout geometry")
    content = None
    for layout in raw["layouts"]:
        n = layout["name"]
        if "Content" in n or "标题和内容" in n or "Title and Content" in n:
            content = layout
            break
    if content is None and len(raw["layouts"]) > 1:
        content = raw["layouts"][1]

    if content is None:
        L.append("Could not identify a content layout.")
    else:
        L.append(f"Primary content layout: `{content['name']}`")
        title_ph = next((s for s in content["shapes"]
                         if s.get("kind") == "placeholder" and s.get("phType") == "title"), None)
        if title_ph:
            L.append("")
            L.append(f"- Title placeholder: x={title_ph.get('x')}, y={title_ph.get('y')}, "
                     f"w={title_ph.get('w')}, h={title_ph.get('h')}")
            L.append(f"- **Recommend:** `space.margin = {title_ph.get('x')}`, "
                     f"`space.titleY = {title_ph.get('y')}`, "
                     f"`space.titleH = {title_ph.get('h')}`")
        else:
            L.append("- No title placeholder found in the content layout.")

    L.append("")
    L.append("---")
    L.append("")
    L.append("Once all calls are resolved, write `theme.json` and delete "
             "this file (and `theme.raw.json`).")

    (out_dir / "extraction_report.md").write_text("\n".join(L))


# ---------- main ------------------------------------------------------------


def main():
    p = argparse.ArgumentParser()
    p.add_argument("--input", "-i", required=True, help="template .pptx")
    p.add_argument("--output-dir", "-o", required=True, help="project directory")
    args = p.parse_args()

    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(args.input) as z:
        master_path, theme_path, master_count = find_active_theme(z)
        theme_info = parse_theme(z, theme_path)

        master_rels = parse_rels(z, rels_path_for(master_path))
        layout_paths = sorted(
            [resolve_rel(master_path, r["target"])
             for r in master_rels if r["type"].endswith("/slideLayout")],
            key=layout_num,
        )
        layouts = [scan_layout(z, lp) for lp in layout_paths]
        sw, sh = slide_dims(z)

    observed = scan_observed_fonts(args.input)

    raw = {
        "source":      str(args.input),
        "master":      master_path,
        "masterCount": master_count,
        "theme":       theme_path,
        "slide":       {"w": sw, "h": sh},
        "colors":      theme_info["colors"],
        "fonts": {
            "declared": theme_info["fonts"],
            "observed": observed,
        },
        "layouts": layouts,
    }

    (out_dir / "theme.raw.json").write_text(
        json.dumps(raw, indent=2, ensure_ascii=False)
    )
    build_report(raw, out_dir)

    print(f"wrote: {out_dir / 'theme.raw.json'}")
    print(f"wrote: {out_dir / 'extraction_report.md'}")


if __name__ == "__main__":
    main()
